# -*- coding: utf-8 -*-
import os, glob, io, base64
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st

# ───────── optional deps/flags ─────────
_SK_OK = True
try:
    from sklearn.linear_model import LinearRegression
except Exception:
    _SK_OK = False

_KALEIDO_OK = True
try:
    import kaleido  # noqa: F401
except Exception:
    _KALEIDO_OK = False

_PDF_OK = True
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas as pdf_canvas
    from reportlab.lib.units import cm
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
except Exception:
    _PDF_OK = False

_AR_SHAPE_OK = True
try:
    import arabic_reshaper
    from bidi.algorithm import get_display
except Exception:
    _AR_SHAPE_OK = False

_PPTX_OK = True
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception:
    _PPTX_OK = False

# ─── شعار من داخل مجلد المشروع (يعمل على أي استضافة) ───
APP_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_DIR = os.path.join(APP_DIR, "assets")          # folder
LOGO_FILE_CANDIDATES = ["logo.png","logo.jpg","logo.jpeg","logo.svg"]

def read_logo_bytes():
    # folder scan first
    if os.path.isdir(LOGO_DIR):
        for name in LOGO_FILE_CANDIDATES:
            p = os.path.join(LOGO_DIR, name)
            if os.path.isfile(p):
                with open(p, "rb") as f:
                    return f.read()
    # fallback: if user put a direct path in LOGO_DIR variable by mistake
    if os.path.isfile(LOGO_DIR):
        with open(LOGO_DIR, "rb") as f:
            return f.read()
    return None

logo_bytes = read_logo_bytes()

# ───────── إعداد عام ─────────
st.set_page_config(page_title="تقرير قسم خدمة العملاء 2025", page_icon="📞", layout="wide")

# ───────── ثيم زجاجي + أحجام موحّدة للـKPI ─────────
DARK_GLASS_CSS = """
<style>
:root{
  --g1: linear-gradient(135deg, rgba(2,132,199,.22), rgba(14,165,233,.10));
  --g2: linear-gradient(135deg, rgba(99,102,241,.22), rgba(14,165,233,.10));
  --g3: linear-gradient(135deg, rgba(16,185,129,.22), rgba(14,165,233,.10));
  --g4: linear-gradient(135deg, rgba(244,114,182,.22), rgba(14,165,233,.10));
  --g5: linear-gradient(135deg, rgba(251,146,60,.22), rgba(14,165,233,.10));
}
html, body, [data-testid="stAppViewContainer"], .block-container { direction: rtl; }
[data-testid="stAppViewContainer"] {
  background: radial-gradient(1200px 800px at 0% 0%, rgba(40,45,60,0.55) 0%, rgba(15,18,25,0.95) 60%),
              radial-gradient(1000px 700px at 100% 0%, rgba(30,35,50,0.55) 0%, rgba(15,18,25,0.95) 60%);
  backdrop-filter: blur(18px);
}
.block-container { padding-top: 2rem; }
.glass {
  background:
    linear-gradient(180deg, rgba(255,255,255,0.09), rgba(255,255,255,0.03)),
    radial-gradient(120% 180% at 100% 0%, rgba(14,165,233,.08), rgba(255,255,255,0));
  border: 1px solid rgba(255,255,255,0.12);
  box-shadow: 0 10px 30px rgba(0,0,0,0.35);
  backdrop-filter: blur(18px); -webkit-backdrop-filter: blur(18px);
  border-radius: 18px; padding: 1.25rem 1.25rem;
}
h1, h2, h3, h4, h5 { color: #ECF2FF !important; letter-spacing: .3px; }
body, p, div, label, span { color: #E3E9F5 !important; }
footer {visibility: hidden;}
button[kind="primary"], .st-emotion-cache-7ym5gk, .st-emotion-cache-19rxjzo {
  background: linear-gradient(135deg, #2b6cb0 0%, #0ea5e9 100%) !important;
  color: white !important; border: none !important; border-radius: 12px !important;
}

/* Header / Logo (أكبر) */
.header-flex { display:flex; align-items:center; justify-content:space-between; padding:1rem 1.25rem; margin-bottom:1rem; }
.header-left { display:flex; gap:.75rem; align-items:center; flex-direction:row-reverse; }
.logo-box {
  width: 84px; height: 84px; border-radius: 18px; overflow:hidden;
  background:linear-gradient(145deg,#0b1222,#0f172a);
  display:flex; align-items:center; justify-content:center;
  box-shadow:0 12px 24px rgba(0,0,0,0.5), inset 0 1px 0 rgba(255,255,255,0.06);
}
.logo-box img { width:100%; height:100%; object-fit:contain; }

/* KPI موحّدة المقاس + تدرّج داخلي */
.kpi {
  position: relative; border-radius: 18px; padding: 1.1rem 1rem;
  text-align: center; min-height: 165px;
  display:flex; flex-direction:column; justify-content:center;
  background:
    linear-gradient(180deg, rgba(255,255,255,0.08), rgba(255,255,255,0.03)),
    radial-gradient(120% 180% at 0% 0%, rgba(255,255,255,0.06), rgba(255,255,255,0));
  border: 1px solid rgba(255,255,255,0.14);
  box-shadow: 0 18px 40px rgba(0,0,0,0.35), inset 0 1px 0 rgba(255,255,255,0.05);
}
.kpi:before {
  content: ""; position: absolute; inset: 0; border-radius: 18px; padding: 1px;
  -webkit-mask: linear-gradient(#000 0 0) content-box, linear-gradient(#000 0 0);
  -webkit-mask-composite: xor; mask-composite: exclude;
}
.kpi.i1:before{background: var(--g1);}
.kpi.i2:before{background: var(--g2);}
.kpi.i3:before{background: var(--g3);}
.kpi.i4:before{background: var(--g4);}
.kpi.i5:before{background: var(--g5);}
.kpi .title { color: #9FB3C8; font-size: .95rem; margin-bottom: .35rem; }
.kpi .value { font-size: 2rem; font-weight: 800; color: #F5FAFF; line-height: 1.15; text-shadow: 0 10px 25px rgba(0,0,0,0.45); }
.kpi .delta-pos { color: #30e3a8; font-weight: 700; }
.kpi .delta-neg { color: #ff7b7b; font-weight: 700; }
.badge {
  display: inline-block; padding: .25rem .55rem; border-radius: 999px; font-size: .85rem;
  background: rgba(255,255,255,0.08); border: 1px solid rgba(255,255,255,0.16); color: #cfe4ff;
}

/* نزّل زر التطبيق شوي */
div[data-testid="stForm"] .stButton button { margin-top: 1.85rem; }
</style>
"""
st.markdown(DARK_GLASS_CSS, unsafe_allow_html=True)

# Header with logo
if logo_bytes:
    b64 = base64.b64encode(logo_bytes).decode("utf-8")
    logo_img_html = f'<div class="logo-box"><img src="data:image/png;base64,{b64}" alt="logo"/></div>'
else:
    logo_img_html = '<div class="logo-box"><span style="font-size:26px;">📞</span></div>'

st.markdown(
    f"""
    <div class="glass header-flex">
      <div class="header-left">
        <div>
          <h2 style="margin:0;">تقرير قسم خدمة العملاء 2025</h2>
          <div style="color:#9FB3C8;margin-top:-4px;font-size:0.95rem;">
            عرض تفاعلي للمؤشرات والرسوم والتنبؤ الذكي
          </div>
        </div>
      </div>
      {logo_img_html}
    </div>
    """,
    unsafe_allow_html=True
)

# ───────── تحميل/تحضير البيانات ─────────
@st.cache_data(show_spinner=True)
def load_and_prepare_data():
    files = glob.glob(os.path.join("data", "*.xlsx"))
    all_dfs = []
    for file in files:
        try:
            df = pd.read_excel(file, sheet_name="الرئيسة", header=11)
        except Exception as e:
            st.error(f"تعذّر قراءة الملف: {os.path.basename(file)} — {e}")
            continue
        df.dropna(how="all", inplace=True)
        agent_name = os.path.basename(file).split("-")[0].strip()
        df["مقدم الخدمة (ملف)"] = agent_name
        all_dfs.append(df)
    if not all_dfs:
        return pd.DataFrame()
    data = pd.concat(all_dfs, ignore_index=True)

    for col in ["اسم العميل ","رقم الجوال ","المنطقة","المدينه ","الشركة","مقدم الخدمة","نوع الخدمة","الخدمه المطلوبه","الملاحظات","الشهر"]:
        if col in data.columns:
            data[col] = data[col].astype(str).str.strip()

    if "المنطقة" in data.columns:
        data["المنطقة"].replace({"المنطقه الشرقيه": "المنطقة الشرقية"}, inplace=True)
    if "نوع الخدمة" in data.columns:
        data["نوع الخدمة"].replace({"nan": None}, inplace=True)
        data["نوع الخدمة"] = data["نوع الخدمة"].astype(str).str.strip()

    # إزالة أرقام تجريبية إن وجدت
    if ("مقدم الخدمة (ملف)" in data.columns) and ("رقم الجوال " in data.columns):
        dummy_phones = [str(num) for num in range(599940931, 599940953)]
        data = data[~((data["مقدم الخدمة (ملف)"] == "Shouq") &
                      (data["رقم الجوال "].astype(str).isin(dummy_phones)))]

    def parse_date(row):
        month_name = str(row.get("الشهر", "")).strip()
        day_val = str(row.get("التاريخ ", "")).strip()
        if not day_val or day_val.lower() == "nan":
            return None
        if "/" in day_val:
            return pd.to_datetime(day_val, dayfirst=True, errors="coerce")
        elif day_val.replace(".", "").isdigit():
            day = int(float(day_val))
            month_map = {"Jan":1,"Feb":2,"Mar":3,"Apr":4,"May":5,"Jun":6,
                         "Jul":7,"Aug":8,"Sep":9,"Oct":10,"Nov":11,"Dec":12}
            m = month_map.get(month_name)
            if m:
                try:
                    return pd.Timestamp(year=2025, month=m, day=day)
                except:
                    return None
        return None

    data["التاريخ/Date"] = data.apply(parse_date, axis=1)
    return data

data = load_and_prepare_data()
if data.empty:
    st.error("لم يتم تحميل أي بيانات. تأكدي من وجود ملفات Excel داخل مجلد data/")
    st.stop()

MONTH_ORDER = ["Aug","Sep","Oct","Nov","Dec","Jan","Feb","Mar","Apr","May","Jun","Jul"]
MONTH_INDEX = {m:i for i,m in enumerate(MONTH_ORDER)}
# ——— توقع الشهر القادم (يُستخدم في بطاقة KPI والرسم لاحقًا) ———
def calc_forecast(df: pd.DataFrame):
    if "الشهر" not in df.columns or df.empty:
        return None, None, None
    month_totals = df["الشهر"].value_counts().reindex(MONTH_ORDER).dropna()
    if len(month_totals) < 1:
        return None, None, None

    mt_index = np.array([MONTH_INDEX[m] for m in month_totals.index]).reshape(-1,1)
    mt_values = month_totals.values.astype(float)

    if _SK_OK and len(mt_values) >= 2:
        model = LinearRegression().fit(mt_index, mt_values)
        next_x = np.array([[mt_index[-1][0] + 1]])
        pred_next = float(model.predict(next_x)[0])
        resid = mt_values - model.predict(mt_index)
        sigma = float(np.std(resid)) if len(resid) > 1 else 0.0
        ci_low = max(0.0, pred_next - 1.96*sigma)
        ci_high = pred_next + 1.96*sigma
        return int(round(pred_next)), int(round(ci_low)), int(round(ci_high))

    if len(mt_values) >= 2:
        growth = mt_values[-1] - mt_values[-2]
        pred_next = int(round(mt_values[-1] + growth))
        return pred_next, None, None

    return None, None, None

# ───────── المرشّحات ─────────
with st.form("filters_form"):
    colf1, colf2, colf3 = st.columns([1.4,1.4,0.9])
    with colf1:
        st.markdown('<div class="glass"><h4 style="margin-top:0;">تصفية حسب الشهر</h4>', unsafe_allow_html=True)
        months_available = sorted(data["الشهر"].dropna().unique().tolist(), key=lambda x: MONTH_INDEX.get(x, 999))
        month_choice = st.selectbox("اختر الشهر", ["الكل"] + months_available, index=0)
        st.markdown('</div>', unsafe_allow_html=True)
    with colf2:
        st.markdown('<div class="glass"><h4 style="margin-top:0;">تصفية حسب مقدم الخدمة</h4>', unsafe_allow_html=True)
        agents_col = "مقدم الخدمة" if "مقدم الخدمة" in data.columns else ("مقدم الخدمة (ملف)" if "مقدم الخدمة (ملف)" in data.columns else None)
        agents_list = sorted(data[agents_col].dropna().unique().tolist()) if agents_col else []
        agent_choice = st.selectbox("اختر مقدم الخدمة", ["الكل"] + agents_list, index=0)
        st.markdown('</div>', unsafe_allow_html=True)
    with colf3:
        st.markdown('<div class="glass"><h4 style="margin-top:0;">تطبيق</h4>', unsafe_allow_html=True)
        st.form_submit_button("تطبيق المرشّحات")
        st.markdown('</div>', unsafe_allow_html=True)

month_final = month_choice if month_choice != "الكل" else None
agent_final = agent_choice if agent_choice != "الكل" else None

filtered = data.copy()
if month_final:
    filtered = filtered[filtered["الشهر"] == month_final]
if agent_final and agents_col:
    filtered = filtered[filtered[agents_col] == agent_final]

# ───────── مؤشرات KPI ─────────
pred_next, pred_ci_low, pred_ci_high = calc_forecast(data)
SERVICE_TYPES = ["استفسار","طلب خدمة","عرض استثماري","طلب تواصل مع أحد منسوبي الشركة","شكوى"]
type_counts = (filtered["نوع الخدمة"].value_counts() if "نوع الخدمة" in filtered.columns else pd.Series(dtype=int))
get_count = lambda k: int(type_counts.get(k, 0))
total_calls = int(len(filtered))

def compute_delta():
    if "الشهر" not in data.columns: return None
    if month_final:
        curr_month = month_final
    else:
        months_present = [m for m in MONTH_ORDER if m in set(filtered["الشهر"].dropna().unique())] if "الشهر" in filtered.columns else []
        if not months_present: return None
        curr_month = months_present[-1]
    curr_idx = MONTH_INDEX.get(curr_month, None)
    if curr_idx is None or curr_idx - 1 < 0: return None
    prev_month = MONTH_ORDER[curr_idx - 1]
    curr_df = data[data["الشهر"] == curr_month]
    prev_df = data[data["الشهر"] == prev_month]
    if agents_col and agent_final:
        curr_df = curr_df[curr_df[agents_col] == agent_final]
        prev_df = prev_df[prev_df[agents_col] == agent_final]
    if curr_df.empty or prev_df.empty: return None
    return int(len(curr_df) - len(prev_df))

kpi_delta = compute_delta()

def top_month_overall():
    if "الشهر" not in data.columns: return None, 0
    df = data.copy()
    if agents_col and agent_final:
        df = df[df[agents_col] == agent_final]
    if df.empty: return None, 0
    s = df["الشهر"].value_counts().sort_values(ascending=False)
    return (s.index[0], int(s.iloc[0])) if len(s) else (None, 0)

def avg_monthly_overall():
    if "الشهر" not in data.columns: return 0.0
    df = data.copy()
    if agents_col and agent_final:
        df = df[df[agents_col] == agent_final]
    if df.empty: return 0.0
    month_tot = df["الشهر"].value_counts()
    return float(month_tot.mean()) if len(month_tot) else 0.0

st.markdown('<div class="glass" style="margin-top:.5rem;">', unsafe_allow_html=True)
st.markdown("### المؤشرات الرئيسية")

# أربعة مؤشرات علوية بنفس المقاس (إجمالي - أعلى شهر - متوسط الشهور - توقع الشهر القادم)
k1,k2,k3,k4 = st.columns([1,1,1,1])

with k1:
    delta_html = ""
    if kpi_delta is not None:
        cls = "delta-pos" if kpi_delta>0 else ("delta-neg" if kpi_delta<0 else "")
        sign = "+" if kpi_delta>0 else ("" if kpi_delta==0 else "")
        delta_html = f'<div class="{cls}" style="margin-top:.3rem;">{sign}{kpi_delta}</div>'
    st.markdown(f"""<div class="kpi i1">
      <div class="title">إجمالي الاتصالات</div>
      <div class="value">{total_calls}</div>
      {delta_html}
      <div style="margin-top:.35rem;color:#9FB3C8;font-size:.9rem;">مقارنةً بالشهر السابق ضمن نفس النطاق</div>
    </div>""", unsafe_allow_html=True)

top_m, top_m_n = top_month_overall()
with k2:
    st.markdown(f"""<div class="kpi i2">
      <div class="title">أعلى شهر</div>
      <div class="value">{top_m or 'غير متاح'}</div>
      <div style="margin-top:.35rem;color:#9FB3C8;font-size:.95rem;">عدد الاتصالات: {int(top_m_n)}</div>
    </div>""", unsafe_allow_html=True)

avg_all = avg_monthly_overall()
with k3:
    st.markdown(f"""<div class="kpi i3">
      <div class="title">متوسط الشهور</div>
      <div class="value">{round(avg_all,1)}</div>
      <div style="margin-top:.35rem;color:#9FB3C8;font-size:.95rem;">متوسط الاتصالات شهريًا على كامل الفترة</div>
    </div>""", unsafe_allow_html=True)

# بطاقة «توقّع الشهر القادم» ثابتة دائمًا
predicted_number = None
ci_low = ci_high = None

with k4:
    if pred_next is not None:
        ci_html = ""
        if pred_ci_low is not None and pred_ci_high is not None:
            ci_html = f'<div class="badge" style="margin-top:.35rem;">نطاق تقريبي: {pred_ci_low} – {pred_ci_high}</div>'
        st.markdown(f"""<div class="kpi i4">
          <div class="title">توقّع الشهر القادم</div>
          <div class="value">{pred_next}</div>
          {ci_html}
        </div>""", unsafe_allow_html=True)
    else:
        st.markdown(f"""<div class="kpi i4">
          <div class="title">توقّع الشهر القادم</div>
          <div class="value">—</div>
          <div class="badge" style="margin-top:.35rem;">البيانات غير كافية</div>
        </div>""", unsafe_allow_html=True)


st.markdown("<div style='height:.5rem;'></div>", unsafe_allow_html=True)

# خمس بطاقات أنواع الخدمات بنفس المقاس
kk1,kk2,kk3,kk4,kk5 = st.columns(5)
for (title, val, cls, col) in [
    ("الاستفسارات", get_count('استفسار'), "i1", kk1),
    ("طلبات الخدمة", get_count('طلب خدمة'), "i2", kk2),
    ("العروض الاستثمارية", get_count('عرض استثماري'), "i3", kk3),
    ("طلبات التواصل", get_count('طلب تواصل مع أحد منسوبي الشركة'), "i4", kk4),
    ("الشكاوى", get_count('شكوى'), "i5", kk5),
]:
    with col:
        st.markdown(f"""<div class="kpi {cls}">
          <div class="title">{title}</div><div class="value">{val}</div>
        </div>""", unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)

# ───────── الرسوم ─────────
figures = {}
st.markdown('<div class="glass" style="margin-top:1rem;">', unsafe_allow_html=True)
st.markdown("### الرسوم التفاعلية")
c1,c2 = st.columns(2)
with c1:
    if "المنطقة" in filtered.columns and not filtered.empty:
        reg_counts = filtered["المنطقة"].value_counts().reset_index()
        reg_counts.columns = ["المنطقة","العدد"]
        fig_reg = px.bar(reg_counts, x="المنطقة", y="العدد", title="توزيع الاتصالات حسب المنطقة", text="العدد")
        fig_reg.update_traces(textposition="outside")
        fig_reg.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                              plot_bgcolor="rgba(255,255,255,0.02)", margin=dict(t=60,b=40,l=20,r=20))
        st.plotly_chart(fig_reg, use_container_width=True)
        figures["توزيع الاتصالات حسب المنطقة"] = fig_reg
    else:
        st.info("لا تتوفر بيانات مناطق للمدى المحدد.")
with c2:
    if "نوع الخدمة" in filtered.columns and not filtered.empty:
        tcounts = filtered["نوع الخدمة"].value_counts()
        fig_type = px.pie(names=tcounts.index, values=tcounts.values, title="نسبة أنواع الاتصالات", hole=0.35)
        fig_type.update_traces(textposition="inside", textinfo="percent+label")
        fig_type.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                               plot_bgcolor="rgba(255,255,255,0.02)", margin=dict(t=60,b=40,l=20,r=20))
        st.plotly_chart(fig_type, use_container_width=True)
        figures["نسبة أنواع الاتصالات"] = fig_type
    else:
        st.info("لا تتوفر بيانات لأنواع الاتصالات للمدى المحدد.")
if "الشركة" in filtered.columns and not filtered.empty:
    comp_counts = filtered["الشركة"].value_counts().reset_index()
    comp_counts.columns = ["الشركة","العدد"]
    fig_comp = px.bar(comp_counts, x="الشركة", y="العدد", title="عدد الاتصالات حسب الشركة", text="العدد")
    fig_comp.update_traces(textposition="outside")
    fig_comp.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                           plot_bgcolor="rgba(255,255,255,0.02)", margin=dict(t=60,b=40,l=20,r=20))
    st.plotly_chart(fig_comp, use_container_width=True)
    figures["عدد الاتصالات حسب الشركة"] = fig_comp
else:
    st.info("لا تتوفر بيانات للشركات للمدى المحدد.")

st.markdown("### توزيع الاتصالات حسب مقدّم الخدمة (على مستوى الشهر الحالي فقط)")
if "الشهر" in data.columns and agents_col:
    if month_final:
        agent_scope_df = data[data["الشهر"] == month_final].copy()
    else:
        months_present = [m for m in MONTH_ORDER if m in set(data["الشهر"].dropna().unique())]
        agent_scope_df = data[data["الشهر"] == months_present[-1]].copy() if months_present else pd.DataFrame()
    if not agent_scope_df.empty:
        agent_counts = agent_scope_df[agents_col].value_counts()
        fig_agents = px.pie(names=agent_counts.index, values=agent_counts.values, hole=0.3,
                            title="نسبة الاتصالات حسب مقدّم الخدمة")
        fig_agents.update_traces(textposition="inside", textinfo="percent+label")
        fig_agents.update_layout(template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                                 plot_bgcolor="rgba(255,255,255,0.02)", margin=dict(t=60,b=40,l=20,r=20))
        st.plotly_chart(fig_agents, use_container_width=True)
        figures["نسبة الاتصالات حسب مقدّم الخدمة"] = fig_agents
st.markdown('</div>', unsafe_allow_html=True)

# ───────── خريطة ─────────
st.markdown('<div class="glass" style="margin-top:1rem;">', unsafe_allow_html=True)
st.markdown("### خريطة الاتصالات حسب المدينة/المنطقة")
CITY_LATLON = {"الرياض": (24.7136, 46.6753),"جدة": (21.4858, 39.1925),"مكة": (21.3891, 39.8579),
               "المدينة": (24.5247, 39.5692),"الدمام": (26.3927, 49.9777),"الخبر": (26.2794, 50.2083),
               "الطائف": (21.2703, 40.4158),"أبها": (18.2465, 42.5117),"حائل": (27.5114, 41.7208),
               "تبوك": (28.3838, 36.5662),"جازان": (16.8892, 42.5700)}
REGION_LATLON = {"المنطقة الشرقية": (26.5, 49.8),"منطقة الرياض": (24.7, 46.7),"منطقة مكة": (21.4, 40.7),
                 "منطقة المدينة": (24.6, 39.6),"منطقة القصيم": (26.3, 43.96),"منطقة تبوك": (28.4, 36.6),
                 "منطقة حائل": (27.5, 41.7),"منطقة جازان": (16.9, 42.6),"منطقة نجران": (17.6, 44.4),
                 "منطقة عسير": (18.2, 42.5),"منطقة الجوف": (29.97, 40.2),"الحدود الشمالية": (30.0, 41.0),
                 "منطقة الباحة": (20.0, 41.45),"منطقة الطائف": (21.27, 40.42)}
def build_map_df(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty: return pd.DataFrame(columns=["label","lat","lon","count"])
    rows = []
    if "المدينه " in df.columns:
        for name, n in df["المدينه "].value_counts().items():
            nm = str(name).strip()
            if nm in CITY_LATLON:
                lat, lon = CITY_LATLON[nm]
                rows.append({"label": nm, "lat": lat, "lon": lon, "count": int(n)})
    if not rows and "المنطقة" in df.columns:
        for name, n in df["المنطقة"].value_counts().items():
            nm = str(name).strip()
            if nm in REGION_LATLON:
                lat, lon = REGION_LATLON[nm]
                rows.append({"label": nm, "lat": lat, "lon": lon, "count": int(n)})
    return pd.DataFrame(rows)
map_df = build_map_df(filtered)
if not map_df.empty:
    fig_map = px.scatter_mapbox(
        map_df, lat="lat", lon="lon", size="count", color="count",
        hover_name="label", hover_data={"lat":False,"lon":False,"count":True},
        size_max=45, zoom=4.2, height=520, title="خريطة توزيع الاتصالات"
    )
    fig_map.update_layout(mapbox_style="carto-darkmatter", paper_bgcolor="rgba(0,0,0,0)",
                          margin=dict(t=60,b=20,l=10,r=10), template="plotly_dark")
    st.plotly_chart(fig_map, use_container_width=True)
    figures["خريطة توزيع الاتصالات"] = fig_map
st.markdown('</div>', unsafe_allow_html=True)

# ───────── جدول + بحث ─────────
st.markdown('<div class="glass" style="margin-top:1rem;">', unsafe_allow_html=True)
st.markdown("### السجل التفصيلي ")
cols_base = ['الشهر','التاريخ ','التاريخ/Date','اسم العميل ','المنطقة','المدينه ','الشركة','نوع الخدمة','مقدم الخدمة','مقدم الخدمة (ملف)']
show_cols = [c for c in cols_base if c in filtered.columns]
q = st.text_input("ابحث داخل الجدول (الاسم/الشركة/المدينة/النوع...)","")
table_df = filtered.copy()
if q.strip() and show_cols:
    ql = q.strip().lower()
    mask = np.zeros(len(table_df), dtype=bool)
    for c in show_cols:
        s = table_df[c].astype(str).str.lower()
        mask |= s.str.contains(ql, na=False)
    table_df = table_df[mask]
if show_cols:
    st.dataframe(table_df[show_cols].reset_index(drop=True), use_container_width=True, height=420)
st.markdown('</div>', unsafe_allow_html=True)

# ───────── تنبؤ + تحديث بطاقة التوقّع ─────────
st.markdown('<div class="glass" style="margin-top:1rem;">', unsafe_allow_html=True)
st.markdown("### التنبؤ بالاتصالات للشهر القادم")
if "الشهر" in data.columns and not data.empty:
    month_totals = data["الشهر"].value_counts().reindex(MONTH_ORDER).dropna()
    if len(month_totals) >= 1:
        mt_index = np.array([MONTH_INDEX[m] for m in month_totals.index]).reshape(-1,1)
        mt_values = month_totals.values.astype(float)
        if _SK_OK and len(mt_values) >= 2:
            model = LinearRegression().fit(mt_index, mt_values)
            next_x = np.array([[mt_index[-1][0] + 1]])
            pred_next = float(model.predict(next_x)[0])
            resid = mt_values - model.predict(mt_index)
            sigma = float(np.std(resid)) if len(resid) > 1 else 0.0
            ci_low = max(0.0, pred_next - 1.96 * sigma)
            ci_high = pred_next + 1.96 * sigma
            line_x = np.arange(mt_index[0][0], next_x[0][0] + 1)
            line_y = model.predict(line_x.reshape(-1,1))
            fig_pred = go.Figure()
            fig_pred.add_trace(go.Scatter(x=[MONTH_ORDER[i] for i in mt_index.flatten()], y=mt_values,
                                          mode="lines+markers", name="فعلي", line=dict(width=3), marker=dict(size=8)))
            fig_pred.add_trace(go.Scatter(x=[MONTH_ORDER[i % 12] for i in line_x], y=line_y,
                                          mode="lines", name="اتجاه", line=dict(dash="dot", width=2)))
            fig_pred.add_trace(go.Scatter(x=[MONTH_ORDER[next_x[0][0] % 12]], y=[pred_next],
                                          mode="markers+text", name="توقع الشهر القادم",
                                          marker=dict(size=12, symbol="diamond"),
                                          text=[f"{int(round(pred_next))}"], textposition="top center"))
            fig_pred.add_trace(go.Scatter(x=[MONTH_ORDER[next_x[0][0] % 12], MONTH_ORDER[next_x[0][0] % 12]],
                                          y=[ci_low, ci_high], mode="lines", name="نطاق تقريبي",
                                          line=dict(color="rgba(164,220,255,.6)", width=8)))
            fig_pred.update_layout(title="منحنى الإجمالي الشهري + التوقع القادم",
                                   template="plotly_dark", paper_bgcolor="rgba(0,0,0,0)",
                                   plot_bgcolor="rgba(255,255,255,0.02)", margin=dict(t=60,b=40,l=20,r=20),
                                   yaxis_title="عدد الاتصالات", xaxis_title="الشهر")
            st.plotly_chart(fig_pred, use_container_width=True)

            # تحديث بطاقة التوقّع أعلى الصفحة (نفس المقاس)
            st.markdown(f"""<div class="kpi i4" style="margin-top:.5rem;">
              <div class="title">توقّع الشهر القادم</div>
              <div class="value">{int(round(pred_next))}</div>
              <div class="badge" style="margin-top:.35rem;">نطاق تقريبي: {int(round(ci_low))} – {int(round(ci_high))}</div>
            </div>""", unsafe_allow_html=True)
        else:
            if len(mt_values) >= 2:
                growth = mt_values[-1] - mt_values[-2]
                pred_next = int(round(mt_values[-1] + growth))
                st.markdown(f"""<div class="kpi i4" style="margin-top:.5rem;">
                <div class="title">توقّع الشهر القادم (تقدير مبسّط)</div>
                <div class="value">{pred_next}</div></div>""", unsafe_allow_html=True)
            else:
                st.info("البيانات غير كافية لحساب التنبؤ.")
st.markdown('</div>', unsafe_allow_html=True)

# ───────── ملخّص ذكي ─────────
st.markdown('<div class="glass" style="margin-top:1rem;">', unsafe_allow_html=True)
st.markdown("### 🤖 ملخّص ذكي للأداء ")
def ar_text(s: str) -> str:
    if not s: return s
    if _AR_SHAPE_OK:
        try:
            return get_display(arabic_reshaper.reshape(s))
        except Exception:
            return s
    return s
def arabic_ai_summary(df: pd.DataFrame) -> str:
    if df.empty: return "لا تتوفر بيانات ضمن النطاق المحدد.."
    top_region = df["المنطقة"].value_counts().index[0] if "المنطقة" in df.columns and not df["المنطقة"].value_counts().empty else None
    top_type   = df["نوع الخدمة"].value_counts().index[0] if "نوع الخدمة" in df.columns and not df["نوع الخدمة"].value_counts().empty else None
    top_comp   = df["الشركة"].value_counts().index[0]   if "الشركة" in df.columns and not df["الشركة"].value_counts().empty else None
    parts = [f"إجمالي الاتصالات في النطاق الحالي بلغ **{len(df)}** اتصالًا."]
    if top_region: parts.append(f"الأكثر نشاطًا: **{top_region}**.")
    if top_type:   parts.append(f"أشيع نوع اتصال: **{top_type}**.")
    if top_comp:   parts.append(f"الشركة الأبرز تواصلًا: **{top_comp}**.")
    if "الشهر" in df.columns:
        by_m = df["الشهر"].value_counts().reindex(MONTH_ORDER).dropna()
        if len(by_m) >= 2:
            delta = int(by_m.iloc[-1] - by_m.iloc[-2])
            trend = "ارتفاع" if delta > 0 else ("انخفاض" if delta < 0 else "ثبات")
            parts.append(f"الاتجاه: **{trend}** بمقدار **{abs(delta)}** مقارنةً بالشهر السابق.")
    return " ".join(parts)
st.write(arabic_ai_summary(filtered))
st.markdown('</div>', unsafe_allow_html=True)

# ───────── Export Helpers ─────────
def fig_to_png_bytes(fig, scale=3.0):
    if not _KALEIDO_OK: return None
    try:
        return fig.to_image(format="png", scale=scale)
    except Exception:
        return None

def register_arabic_font() -> str | None:
    if not _PDF_OK: return None
    candidates = [("NotoNaskhArabic-Regular", "fonts/NotoNaskhArabic-Regular.ttf"),
                  ("Amiri-Regular", "fonts/Amiri-Regular.ttf"),
                  ("Cairo-Regular", "fonts/Cairo-Regular.ttf")]
    for fam, path in candidates:
        if os.path.isfile(path):
            try: pdfmetrics.registerFont(TTFont(fam, path)); return fam
            except Exception: pass
    return None

# ───────── Export: PNG / PDF / PPTX ─────────
st.markdown('<div class="glass" style="margin-top:1rem;">', unsafe_allow_html=True)
st.markdown("### ⬇️ التصدير (PNG / PDF / PPTX)")
col_dl1, col_dl2, col_dl3 = st.columns(3)

# PNG
with col_dl1:
    st.markdown("**تصدير الرسوم كصور PNG**")
    if not figures:
        st.info("لا توجد رسوم جاهزة للتصدير ضمن النطاق الحالي.")
    for name, fig in figures.items():
        png_buf = fig_to_png_bytes(fig, scale=3.0)
        if png_buf:
            st.download_button(f"تحميل {name} كصورة PNG", data=png_buf,
                               file_name=f"{name}.png", mime="image/png",
                               use_container_width=True)
        else:
            st.caption(f"يتطلب kaleido: pip install kaleido — ({name})")

# PDF (منظم + شعار أعلى اليمين)
with col_dl2:
    st.markdown("**تقرير PDF منظم (يدعم العربية + شعار ثابت أعلى اليمين)**")
    can_pdf = _PDF_OK and _KALEIDO_OK and (len(figures) > 0)
    if not _PDF_OK: st.caption("يتطلب reportlab")
    if not _KALEIDO_OK: st.caption("يتطلب kaleido")
    if can_pdf:
        ar_font = register_arabic_font()
        pdf_buffer = io.BytesIO()
        c = pdf_canvas.Canvas(pdf_buffer, pagesize=A4)
        W, H = A4
        M_L, M_R, M_T, M_B = 2*cm, 2*cm, 2*cm, 2*cm

        def draw_header(title_text: str):
            if logo_bytes:
                try:
                    img = ImageReader(io.BytesIO(logo_bytes))
                    logo_w = 2.8*cm; logo_h = logo_w*0.9
                    c.drawImage(img, W - M_R - logo_w, H - M_T - logo_h,
                                width=logo_w, height=logo_h, preserveAspectRatio=True, mask='auto')
                except Exception:
                    pass
            if ar_font and _AR_SHAPE_OK:
                c.setFont(ar_font, 16)
                c.drawRightString(W - M_R - (0.3*cm if logo_bytes else 0), H - M_T - 0.3*cm,
                                  get_display(arabic_reshaper.reshape(title_text)))
            else:
                c.setFont("Helvetica-Bold", 16)
                c.drawString(M_L, H - M_T - 0.3*cm, title_text)

        def draw_rtl(text: str, x_right: float, y: float, font_name: str, font_size: int):
            c.setFont(font_name, font_size)
            if ar_font and _AR_SHAPE_OK:
                c.drawRightString(x_right, y, get_display(arabic_reshaper.reshape(text)))
            else:
                c.drawString(M_L, y, text)

        # صفحة موجزة
        draw_header("لوحة مركز الاتصال — تقرير مختصر")
        y = H - M_T - 2.0*cm
        line_font = ar_font if ar_font else "Helvetica"
        for ln in [
            f"إجمالي الاتصالات: {total_calls}",
            (f"الفارق عن الشهر السابق (ضمن النطاق): {('+' if (kpi_delta or 0)>0 else '')}{kpi_delta}") if kpi_delta is not None else "",
        ]:
            if ln:
                draw_rtl(ln, W - M_R, y, line_font, 11); y -= 0.6*cm

        summary = arabic_ai_summary(filtered)
        for piece in summary.split(". "):
            draw_rtl(piece if piece.endswith(".") else piece + ".", W - M_R, y, line_font, 10)
            y -= 0.5*cm
            if y < M_B + 4*cm:
                c.showPage(); draw_header("متابعة التقرير"); y = H - M_T - 2.0*cm

        # شبكة 2×2
        slot_w = (W - M_L - M_R - 1.0*cm) / 2.0
        slot_h = (H - M_T - M_B - 3.0*cm) / 2.0
        Xs = [M_L, M_L + slot_w + 1.0*cm]
        Ys = [H - M_T - slot_h - 1.0*cm, M_B + 1.5*cm]

        c.showPage(); draw_header("رسوم مختارة")
        idx = 0
        items = list(figures.items())
        for name, fig in items:
            img_bytes = fig_to_png_bytes(fig, scale=3.0)
            if not img_bytes: continue
            img = ImageReader(io.BytesIO(img_bytes))
            col = idx % 2
            row = 0 if (idx // 2) % 2 == 0 else 1
            x = Xs[col]; y_img = Ys[row]
            c.drawImage(img, x, y_img, width=slot_w, height=slot_h,
                        preserveAspectRatio=True, anchor='sw', mask='auto')
            c.setFont(line_font, 9)
            cap = get_display(arabic_reshaper.reshape(name)) if (ar_font and _AR_SHAPE_OK) else name
            if ar_font and _AR_SHAPE_OK:
                c.drawRightString(x + slot_w, y_img - 0.25*cm, cap)
            else:
                c.drawString(x, y_img - 0.25*cm, name)

            if (idx % 4) == 3 and idx != len(items)-1:
                c.showPage(); draw_header("متابعة الرسوم")
            idx += 1

        c.save(); pdf_buffer.seek(0)
        st.download_button("تحميل تقرير PDF (منظم + شعار ثابت)", data=pdf_buffer,
                           file_name="تقرير_مركز_الاتصال.pdf", mime="application/pdf",
                           use_container_width=True)

# PPTX
with col_dl3:
    st.markdown("**PowerPoint (شعار ثابت أعلى اليمين + شريحة لكل رسم)**")
    if _PPTX_OK and _KALEIDO_OK and (len(figures) > 0):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.5), Inches(1)).text_frame
        tf.text = "لوحة مركز الاتصال — رسوم مختارة"
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        if logo_bytes:
            slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(9.1), Inches(0.3), width=Inches(1.0))

        for name, fig in figures.items():
            img_bytes = fig_to_png_bytes(fig, scale=3.0)
            if not img_bytes: continue
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(0.7)).text_frame
            tx.text = name; tx.paragraphs[0].font.size = Pt(22); tx.paragraphs[0].alignment = PP_ALIGN.RIGHT
            if logo_bytes:
                slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(9.1), Inches(0.25), width=Inches(1.0))
            slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(1.0), width=Inches(9))

        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        st.download_button("تحميل PowerPoint", data=buf,
                           file_name="CallCenter_Dashboard_Charts.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           use_container_width=True)
    else:
        if not _PPTX_OK: st.caption("يتطلب python-pptx")
        if not _KALEIDO_OK: st.caption("يتطلب kaleido")
        if len(figures) == 0: st.caption("لا توجد رسوم متاحة.")
st.markdown('</div>', unsafe_allow_html=True)

